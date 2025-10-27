from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slides = [
    {
        'title': 'Smart Surveillance System',
        'bullets': ['Real-time detection: weapons, crowds, suspicious behavior', 'Automatic alerts: email + audible beeps', 'Auth & GUI for operators'],
        'notes': "Start with one-line value prop. Say you built an integrated demo that combines object detection, emotion analysis and optional violence classification."
    },
    {
        'title': 'Problem & Motivation',
        'bullets': ['Manual monitoring is slow and error-prone', 'Need automated alerts to improve response time', 'Low-cost, deployable prototype for safety'],
        'notes': 'Explain the high-level problem and why automation matters for safety and response.'
    },
    {
        'title': 'System Architecture',
        'bullets': ['Camera capture → Preprocess → Models → Alerts → GUI/DB', 'Async inference keeps UI responsive', 'Persistent boxes and evidence capture'],
        'notes': 'Walk through the flow. Emphasize async inference and evidence email.'
    },
    {
        'title': 'Models & Roles',
        'bullets': ['YOLOv8: weapon + person detection (per-frame)', 'FER: face emotion estimation (support signal)', 'SlowFast (PyTorchVideo): clip-level violence detection (optional)'],
        'notes': 'Briefly explain each model and why it fits its role (spatial vs temporal).'
    },
    {
        'title': 'Key Implementation Decisions',
        'bullets': ['Detection runs asynchronously (ThreadPoolExecutor)', 'Lower capture resolution for performance', 'Env-var email config and DB logging'],
        'notes': 'Explain design reasons: responsiveness, false positive mitigation, secure config.'
    },
    {
        'title': 'Demo & Metrics',
        'bullets': ['Live demo: show detection + email alert', 'Metrics: latency (ms), detection frequency, false positive rate'],
        'notes': 'Show a short clip and email; present any latency numbers or explain CPU/GPU differences.'
    },
    {
        'title': 'Limitations & Mitigations',
        'bullets': ['Violence model is heavy — best on GPU', 'YOLO false positives on metal; use thresholds & tracking', 'FER unstable on low-res faces'],
        'notes': 'Be honest about limits and describe how you mitigate them.'
    },
    {
        'title': 'Roadmap & Ask',
        'bullets': ['GPU access & benchmarking', 'Model optimization (ONNX/TensorRT)', 'User studies and multi-camera scaling'],
        'notes': 'Close with a clear ask: resources, time, and the next steps. Invite questions.'
    }
]

for s in slides:
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    title.text = s['title']
    tf = body.text_frame
    tf.clear()
    for i, b in enumerate(s['bullets']):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = b
        else:
            p = tf.add_paragraph()
            p.text = b
        p.level = 0
        p.font.size = Pt(28)
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = s['notes']

out_path = 'docs/Review_Presentation.pptx'
prs.save(out_path)
print('Saved', out_path)
